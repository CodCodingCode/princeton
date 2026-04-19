#!/usr/bin/env python3
"""Build the Onkos pitch deck (onkos_pitch.pptx) at repo root.

Biotech-clean, 16:9, navy-on-white. Ten slides mapped to the HackPrinceton
Regeneron judging rubric: readiness, utility, design, relevance, packaging.
Run: python3 scripts/build_deck.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Palette (matches frontend/tailwind.config.ts)
NAVY_900 = RGBColor(0x06, 0x16, 0x30)
NAVY_700 = RGBColor(0x0B, 0x25, 0x45)
NAVY_500 = RGBColor(0x1E, 0x3A, 0x8A)
NAVY_50  = RGBColor(0xEE, 0xF2, 0xF8)
INK_0    = RGBColor(0x0A, 0x0A, 0x0A)
INK_400  = RGBColor(0x6B, 0x6B, 0x6B)
INK_600  = RGBColor(0xA3, 0xA3, 0xA3)
INK_800  = RGBColor(0xE7, 0xE7, 0xE7)
INK_900  = RGBColor(0xF4, 0xF4, 0xF4)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DOT_LIVE = RGBColor(0x10, 0xB9, 0x81)

FONT       = "Calibri"
FONT_LIGHT = "Calibri Light"

TOTAL_SLIDES = 10
TRACK_LINE   = "HackPrinceton 2026  ·  Regeneron Track"


# Basic helpers
def new_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def textbox(slide, x, y, w, h, text, *, font=FONT, size=14, bold=False,
            color=INK_0, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            tracking=0, line_spacing=1.2, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if tracking:
        run._r.get_or_add_rPr().set("spc", str(tracking))
    return tb


def rect(slide, x, y, w, h, fill=NAVY_700, line=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_width is not None:
            shp.line.width = line_width
    return shp


def oval(slide, x, y, d, fill=NAVY_700):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def hline(slide, x, y, w, color=INK_600, width=Pt(0.5)):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y)
    c.line.color.rgb = color
    c.line.width = width
    return c


def arrow(slide, x, y, w, color=INK_600, width=Pt(0.75)):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y)
    c.line.color.rgb = color
    c.line.width = width
    ln = c.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("h", "sm")
    return c


def eyebrow(slide, x, y, num, label):
    d = Inches(0.09)
    oval(slide, x, y + Inches(0.06), d, fill=NAVY_700)
    textbox(slide, x + d + Inches(0.12), y, Inches(6), Inches(0.22),
            f"{num:02d}  ·  {label.upper()}",
            font=FONT, size=9, bold=True, color=NAVY_700,
            anchor=MSO_ANCHOR.MIDDLE, tracking=200)


def chrome(slide, slide_num, section):
    if slide_num > 0:
        eyebrow(slide, Inches(0.6), Inches(0.5), slide_num, section)
    textbox(slide, Inches(8.5), Inches(0.5), Inches(4.233), Inches(0.22),
            "ONKOS", font=FONT, size=9, bold=True, color=INK_400,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, tracking=300)
    hline(slide, Inches(0.6), Inches(7.1), Inches(12.133), color=INK_800)
    textbox(slide, Inches(0.6), Inches(7.18), Inches(6), Inches(0.2),
            f"{slide_num:02d} / {TOTAL_SLIDES:02d}" if slide_num > 0 else "COVER",
            font=FONT, size=8, color=INK_400, tracking=150)
    textbox(slide, Inches(6.7), Inches(7.18), Inches(6.033), Inches(0.2),
            TRACK_LINE, font=FONT, size=8, color=INK_400,
            align=PP_ALIGN.RIGHT, tracking=150)


def headline(slide, text, size=30, top=None, color=INK_0):
    y = top if top is not None else Inches(1.05)
    return textbox(slide, Inches(0.6), y, Inches(12.133), Inches(1.2),
                   text, font=FONT_LIGHT, size=size, color=color,
                   line_spacing=1.1)


def subhead(slide, text, top):
    return textbox(slide, Inches(0.6), top, Inches(12.133), Inches(0.8),
                   text, font=FONT, size=15, color=INK_400,
                   line_spacing=1.35)


# SLIDE 1 — Cover
def slide_cover(prs):
    slide = blank_slide(prs)
    rect(slide, Inches(0), Inches(0), Inches(0.22), Inches(7.5), fill=NAVY_700)
    textbox(slide, Inches(0.9), Inches(0.5), Inches(6), Inches(0.3),
            "ONKOS", font=FONT, size=10, bold=True, color=NAVY_700,
            tracking=400)
    textbox(slide, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.6),
            "Cursor",
            font=FONT_LIGHT, size=72, color=INK_0, line_spacing=1.0)
    textbox(slide, Inches(0.9), Inches(3.55), Inches(11.5), Inches(1.2),
            "for Oncologists.",
            font=FONT_LIGHT, size=72, color=NAVY_700, line_spacing=1.0)
    textbox(slide, Inches(0.9), Inches(4.95), Inches(11.5), Inches(0.7),
            "Upload a patient's records. Get back a cited treatment plan, "
            "matching trials, and a clinician report in minutes.",
            font=FONT_LIGHT, size=16, color=INK_400, line_spacing=1.4)
    hline(slide, Inches(0.9), Inches(6.55), Inches(11.5), color=INK_800)
    textbox(slide, Inches(0.9), Inches(6.65), Inches(6), Inches(0.3),
            "HACKPRINCETON 2026  ·  REGENERON TRACK",
            font=FONT, size=9, bold=True, color=NAVY_700, tracking=250)
    textbox(slide, Inches(7), Inches(6.65), Inches(5.4), Inches(0.3),
            "TEAM ONKOS",
            font=FONT, size=9, bold=True, color=INK_400,
            align=PP_ALIGN.RIGHT, tracking=250)


# SLIDE 2 — Problem
def slide_problem(prs):
    slide = blank_slide(prs)
    chrome(slide, 1, "The problem")
    headline(slide, "Matching a patient to the right trial\nstill takes days, not minutes.",
             size=34)
    subhead(slide, "Oncologists piece together PDFs, scans, and handwritten notes by hand. "
                   "The matching criteria sit buried inside hundreds of trial protocols.",
            top=Inches(2.75))

    cards = [
        ("86%",          "of clinical trials miss their original enrollment timelines."),
        ("200+",         "pages of records behind a single new patient's chart."),
        ("Days",         "to hand-screen one patient against a trial shortlist."),
    ]
    card_w = Inches(3.9)
    gap = Inches(0.25)
    y_card = Inches(4.35)
    total_w = card_w * 3 + gap * 2
    x0 = int((Inches(13.333) - total_w) / 2)
    for i, (stat, cap) in enumerate(cards):
        x = x0 + (card_w + gap) * i
        rect(slide, x, y_card, card_w, Inches(2.1), fill=INK_900)
        rect(slide, x, y_card, Inches(0.06), Inches(2.1), fill=NAVY_700)
        textbox(slide, x + Inches(0.35), y_card + Inches(0.3),
                card_w - Inches(0.6), Inches(0.9),
                stat, font=FONT_LIGHT, size=40, color=NAVY_700,
                line_spacing=1.0)
        textbox(slide, x + Inches(0.35), y_card + Inches(1.2),
                card_w - Inches(0.6), Inches(0.8),
                cap, font=FONT, size=12, color=INK_400, line_spacing=1.35)


# SLIDE 3 — Security & privacy
def slide_security(prs):
    slide = blank_slide(prs)
    chrome(slide, 2, "Security & privacy")
    headline(slide, "HIPAA-aligned by design.", size=32)
    subhead(slide,
            "No real patient data enters the system. Logs are scrubbed before they "
            "touch disk, endpoints can be token-gated, and case state never leaves memory.",
            top=Inches(2.2))

    cards = [
        ("01", "SYNTHETIC DATA", [
            "No real patient records, ever.",
            "Every demo case is generated from scratch.",
            "No PHI, no IRB, no burden.",
        ]),
        ("02", "REDACTED LOGS", [
            "Regex scrubber runs on every log line.",
            "Catches SSN, DOB, phone, email, MRN.",
            "On by default (NEOVAX_LOG_REDACTION).",
        ]),
        ("03", "TOKEN-GATED API", [
            "Every /api/* can require a bearer token.",
            "Constant-time comparison (secrets.compare_digest).",
            "Flip on with NEOVAX_API_TOKEN.",
        ]),
        ("04", "EPHEMERAL CASES", [
            "Cases live in memory for one session.",
            "Nothing persists to disk.",
            "A server restart drops every case.",
        ]),
    ]

    y = Inches(3.4)
    h = Inches(2.8)
    gap = Inches(0.22)
    n = 4
    avail = Inches(12.133) - gap * (n - 1)
    w = int(avail / n)
    x = Inches(0.6)
    for num, title, items in cards:
        rect(slide, x, y, w, h, fill=WHITE, line=INK_800, line_width=Pt(0.75))
        rect(slide, x, y, Inches(0.06), h, fill=NAVY_700)
        textbox(slide, x + Inches(0.3), y + Inches(0.22),
                w - Inches(0.5), Inches(0.35),
                num, font=FONT_LIGHT, size=20, color=NAVY_700, tracking=100)
        textbox(slide, x + Inches(0.3), y + Inches(0.58),
                w - Inches(0.5), Inches(0.3),
                title, font=FONT, size=10, bold=True, color=INK_0, tracking=250)
        hline(slide, x + Inches(0.3), y + Inches(0.92),
              w - Inches(0.5), color=INK_800)
        by = y + Inches(1.05)
        for b in items:
            oval(slide, x + Inches(0.3), by + Inches(0.08), Inches(0.07),
                 fill=NAVY_500)
            textbox(slide, x + Inches(0.48), by, w - Inches(0.7),
                    Inches(0.5), b, font=FONT, size=10.5, color=INK_0,
                    line_spacing=1.3)
            by += Inches(0.5)
        x += w + gap

    # Honest scope note along the bottom — don't let judges think we're claiming
    # full HIPAA certification, since we aren't.
    y2 = Inches(6.3)
    rect(slide, Inches(0.6), y2, Inches(12.133), Inches(0.7), fill=NAVY_50)
    textbox(slide, Inches(0.85), y2 + Inches(0.1), Inches(5), Inches(0.25),
            "ALIGNED, NOT CERTIFIED",
            font=FONT, size=9, bold=True, color=NAVY_700, tracking=300)
    textbox(slide, Inches(0.85), y2 + Inches(0.34), Inches(11.6), Inches(0.3),
            "Full HIPAA compliance needs a BAA with every vendor (K2-Think, Kimi, HeyGen, Google Maps). "
            "This is the foundation, not the finish line.",
            font=FONT, size=10, color=INK_400, line_spacing=1.3)


# SLIDE 4 — Architecture
def slide_architecture(prs):
    slide = blank_slide(prs)
    chrome(slide, 3, "Architecture")
    headline(slide, "How the pieces fit together.", size=32)
    subhead(slide, "A Next.js frontend talks to a FastAPI backend over server-sent events. "
                   "Model output is pushed to the UI as soon as it arrives.",
            top=Inches(2.2))

    y = Inches(3.4)
    col_w = Inches(3.9)
    gap = Inches(0.3)
    x0 = int((Inches(13.333) - (col_w * 3 + gap * 2)) / 2)
    columns = [
        ("FRONTEND", "Next.js 15  ·  Tailwind", [
            "Upload flow and live cockpit",
            "Clinician tabs: Overview, Plan,\nTrials, Documents, Clinical",
            "Patient sidebar with spoken avatar",
            "Tab state lives in the URL",
        ]),
        ("BACKEND", "FastAPI  ·  SSE  ·  Python 3.11", [
            "Orchestrator walks the pipeline",
            "Case store and event bus in memory",
            "PDF report generator",
            "Google Maps geocoding",
        ]),
        ("MODELS & DATA", "K2-Think  ·  Kimi  ·  HeyGen", [
            "K2-Think: reasons through\nguidelines, shows its work",
            "Kimi: handles chat with tools",
            "HeyGen: gives the avatar a voice",
            "PubMed index for citations",
            "Regeneron trial registry",
        ]),
    ]
    card_h = Inches(3.45)
    for i, (title, stack, bullets) in enumerate(columns):
        x = x0 + (col_w + gap) * i
        rect(slide, x, y, col_w, card_h, fill=WHITE, line=INK_800, line_width=Pt(0.75))
        rect(slide, x, y, col_w, Inches(0.06), fill=NAVY_700)
        textbox(slide, x + Inches(0.35), y + Inches(0.25),
                col_w - Inches(0.7), Inches(0.3),
                title, font=FONT, size=10, bold=True, color=NAVY_700, tracking=250)
        textbox(slide, x + Inches(0.35), y + Inches(0.55),
                col_w - Inches(0.7), Inches(0.3),
                stack, font=FONT, size=11, color=INK_400, tracking=50)
        hline(slide, x + Inches(0.35), y + Inches(0.9),
              col_w - Inches(0.7), color=INK_800)
        by = y + Inches(1.05)
        for b in bullets:
            oval(slide, x + Inches(0.35), by + Inches(0.08), Inches(0.07),
                 fill=NAVY_500)
            textbox(slide, x + Inches(0.55), by, col_w - Inches(0.9),
                    Inches(0.55), b, font=FONT, size=11, color=INK_0,
                    line_spacing=1.3)
            lines = b.count("\n") + 1
            by += Inches(0.32 + 0.18 * (lines - 1))


# SLIDE 5 — Synthetic data (IMAGE PLACEHOLDERS)
def slide_synthetic_data(prs):
    slide = blank_slide(prs)
    chrome(slide, 4, "Synthetic data")
    headline(slide, "Built on realistic, made-up case bundles.", size=32)
    subhead(slide,
            "Every demo case is synthetic: pathology PDFs, faxed referrals, lab charts, "
            "registration sheets. The same shape and mess a real intake looks like, with no real patient data.",
            top=Inches(2.2))

    labels = [
        ("FAX REFERRAL",   "Scanned outside referral\nwith handwriting and stamps."),
        ("LAB CHART",      "CBC and CMP values\ntracked over time."),
        ("PATHOLOGY PDF",  "Digital biopsy report\nwith IHC and molecular panel."),
    ]
    card_w = Inches(3.9)
    gap = Inches(0.3)
    y = Inches(3.45)
    card_h = Inches(3.35)
    x0 = int((Inches(13.333) - (card_w * 3 + gap * 2)) / 2)
    for i, (title, caption) in enumerate(labels):
        x = x0 + (card_w + gap) * i
        rect(slide, x, y, card_w, card_h, fill=WHITE,
             line=INK_800, line_width=Pt(0.75))
        drop_x = x + Inches(0.2)
        drop_y = y + Inches(0.2)
        drop_w = card_w - Inches(0.4)
        drop_h = Inches(2.15)
        rect(slide, drop_x, drop_y, drop_w, drop_h, fill=INK_900,
             line=INK_600, line_width=Pt(0.5))
        textbox(slide, drop_x, int(drop_y + drop_h / 2 - Inches(0.25)),
                drop_w, Inches(0.3),
                "DROP IMAGE HERE", font=FONT, size=9, bold=True,
                color=INK_400, align=PP_ALIGN.CENTER, tracking=350)
        textbox(slide, drop_x, int(drop_y + drop_h / 2 + Inches(0.05)),
                drop_w, Inches(0.3),
                "(replace this rectangle with the screenshot)",
                font=FONT, size=9, color=INK_600,
                align=PP_ALIGN.CENTER, italic=True)
        textbox(slide, x + Inches(0.25), y + Inches(2.6),
                card_w - Inches(0.5), Inches(0.28),
                title, font=FONT, size=10, bold=True, color=NAVY_700, tracking=250)
        textbox(slide, x + Inches(0.25), y + Inches(2.95),
                card_w - Inches(0.5), Inches(0.5),
                caption, font=FONT, size=11, color=INK_400, line_spacing=1.35)


# SLIDE 6 — Regeneron relevance
def slide_regeneron(prs):
    slide = blank_slide(prs)
    chrome(slide, 5, "Regeneron relevance")
    headline(slide, "Built around the Regeneron oncology pipeline.",
             size=30)
    subhead(slide,
            "Trial matching runs on real eligibility rules from real Regeneron-sponsored "
            "studies. Not a keyword search of ClinicalTrials.gov.",
            top=Inches(2.2))

    y = Inches(3.4)
    tiles = [
        ("29",     "Regeneron trials\ncovered."),
        ("12",     "Biomarker rules\nbuilt in."),
        ("16",     "Eligibility checks\nper trial."),
        ("3-way",  "Eligible, needs data,\nor ineligible verdict."),
    ]
    tile_w = Inches(3.0)
    gap = Inches(0.15)
    x0 = int((Inches(13.333) - (tile_w * 4 + gap * 3)) / 2)
    for i, (stat, cap) in enumerate(tiles):
        x = x0 + (tile_w + gap) * i
        rect(slide, x, y, tile_w, Inches(1.5), fill=NAVY_50)
        textbox(slide, x + Inches(0.3), y + Inches(0.22),
                tile_w - Inches(0.6), Inches(0.7),
                stat, font=FONT_LIGHT, size=32, color=NAVY_700, line_spacing=1.0)
        textbox(slide, x + Inches(0.3), y + Inches(0.85),
                tile_w - Inches(0.6), Inches(0.6),
                cap, font=FONT, size=11, color=INK_400, line_spacing=1.3)

    y2 = Inches(5.2)
    textbox(slide, Inches(0.6), y2, Inches(12.133), Inches(0.3),
            "KEY REGENERON ASSETS COVERED",
            font=FONT, size=9, bold=True, color=NAVY_700, tracking=250)
    hline(slide, Inches(0.6), y2 + Inches(0.3), Inches(12.133), color=INK_800)
    programs = [
        ("Fianlimab + cemiplimab",  "LAG-3 with PD-1 in advanced melanoma."),
        ("Cemiplimab monotherapy",  "PD-1 across solid tumors."),
        ("Odronextamab",            "CD20×CD3 bispecific for B-cell lymphoma."),
        ("Linvoseltamab",           "BCMA×CD3 bispecific for multiple myeloma."),
    ]
    col_w = Inches(3.0)
    gap2 = Inches(0.05)
    yr = y2 + Inches(0.45)
    x0p = int((Inches(13.333) - (col_w * 4 + gap2 * 3)) / 2)
    for i, (name, desc) in enumerate(programs):
        x = x0p + (col_w + gap2) * i
        textbox(slide, x + Inches(0.1), yr, col_w - Inches(0.2), Inches(0.3),
                name, font=FONT, size=12, bold=True, color=INK_0)
        textbox(slide, x + Inches(0.1), yr + Inches(0.35),
                col_w - Inches(0.2), Inches(0.8),
                desc, font=FONT, size=10.5, color=INK_400, line_spacing=1.3)


# SLIDE 7 — Reasoning engine
def slide_reasoning(prs):
    slide = blank_slide(prs)
    chrome(slide, 6, "How it reasons")
    headline(slide, "The reasoning model drives it.", size=32)
    subhead(slide,
            "The NCCN treatment pathway is actually walked by a reasoning model, step by step. "
            "You watch it think in real time.",
            top=Inches(2.2))

    x_left = Inches(0.6)
    y = Inches(3.4)
    bullets = [
        ("Live reasoning stream",
         "K2-Think thinks out loud at each guideline step. The cockpit renders it as it happens."),
        ("Grounded in real papers",
         "Each step pulls relevant PubMed citations, filtered to the patient's cancer type."),
        ("Works across cancer types",
         "The walker picks the right branch based on the cancer type the extractor found. Nothing is hard-coded."),
        ("Follow-up chat",
         "A second model handles questions after the run. Its tools only guide the UI; they never change the case."),
        ("The avatar",
         "Kimi writes the answer. HeyGen gives it a voice. The avatar doesn't think on its own."),
    ]
    left_w = Inches(7.5)
    row_h = Inches(0.72)
    for i, (t, d) in enumerate(bullets):
        by = y + row_h * i
        oval(slide, x_left, by + Inches(0.18), Inches(0.1), fill=NAVY_700)
        textbox(slide, x_left + Inches(0.25), by,
                left_w - Inches(0.3), Inches(0.3),
                t, font=FONT, size=13, bold=True, color=INK_0)
        textbox(slide, x_left + Inches(0.25), by + Inches(0.3),
                left_w - Inches(0.3), Inches(0.42),
                d, font=FONT, size=10.5, color=INK_400, line_spacing=1.3)

    px = Inches(8.7)
    pw = Inches(4.033)
    ph = Inches(3.55)
    py = Inches(3.4)
    rect(slide, px, py, pw, ph, fill=NAVY_900)
    rect(slide, px, py, pw, Inches(0.4), fill=NAVY_700)
    oval(slide, px + Inches(0.2), py + Inches(0.14), Inches(0.12), fill=DOT_LIVE)
    textbox(slide, px + Inches(0.45), py + Inches(0.08), pw - Inches(0.6),
            Inches(0.3), "K2-THINK  ·  LIVE", font=FONT, size=9, bold=True,
            color=WHITE, tracking=250, anchor=MSO_ANCHOR.MIDDLE)
    think = (
        "thinking... patient presents with Stage IIIC cutaneous melanoma, "
        "BRAF V600E positive, ECOG 1. primary decision: adjuvant vs. "
        "neoadjuvant systemic therapy.\n\n"
        "per NCCN MEL-C: for resectable stage III with positive "
        "sentinel node, consider neoadjuvant pembrolizumab or nivolumab "
        "plus relatlimab (RELATIVITY-098).\n\n"
        "flag: patient may also be eligible for a LAG-3 combo. "
        "cross-check Regeneron fianlimab + cemiplimab (NCT05352672)..."
    )
    textbox(slide, px + Inches(0.3), py + Inches(0.55), pw - Inches(0.6),
            ph - Inches(0.7), think, font=FONT, size=10,
            color=RGBColor(0xD7, 0xDF, 0xF0), line_spacing=1.35, italic=True)


# SLIDE 8 — Design: dual-audience cockpit
def slide_design(prs):
    slide = blank_slide(prs)
    chrome(slide, 7, "Design")
    headline(slide, "One case. Two audiences.", size=32)
    subhead(slide,
            "The same case data powers two views. The avatar can send the user "
            "to the right tab when it needs to.",
            top=Inches(2.2))

    y = Inches(3.4)
    h = Inches(3.4)
    gap = Inches(0.3)
    w = int((Inches(12.133) - gap) / 2)
    x_left = Inches(0.6)
    x_right = x_left + w + gap

    def mock_panel(x, kind, title, tabs, bullets):
        rect(slide, x, y, w, h, fill=WHITE, line=INK_800, line_width=Pt(0.75))
        rect(slide, x, y, w, Inches(0.55), fill=NAVY_50)
        textbox(slide, x + Inches(0.35), y + Inches(0.1), w - Inches(0.7),
                Inches(0.25), kind, font=FONT, size=9, bold=True,
                color=NAVY_700, tracking=300)
        textbox(slide, x + Inches(0.35), y + Inches(0.3), w - Inches(0.7),
                Inches(0.25), title, font=FONT, size=13, bold=True, color=INK_0)
        # tab strip, sized to fit inside card
        inside_w = w - Inches(0.7)
        tgap = Inches(0.06)
        tw = int((inside_w - tgap * (len(tabs) - 1)) / len(tabs))
        tby = y + Inches(0.75)
        tx = x + Inches(0.35)
        for i, tab in enumerate(tabs):
            active = (i == 0)
            if active:
                rect(slide, tx, tby, tw, Inches(0.3), fill=NAVY_700)
                textbox(slide, tx, tby + Inches(0.05), tw, Inches(0.22),
                        tab, font=FONT, size=9, bold=True, color=WHITE,
                        align=PP_ALIGN.CENTER, tracking=150)
            else:
                rect(slide, tx, tby, tw, Inches(0.3), fill=INK_900)
                textbox(slide, tx, tby + Inches(0.05), tw, Inches(0.22),
                        tab, font=FONT, size=9, color=INK_400,
                        align=PP_ALIGN.CENTER, tracking=150)
            tx += tw + tgap
        by = y + Inches(1.3)
        for b in bullets:
            oval(slide, x + Inches(0.4), by + Inches(0.09), Inches(0.07),
                 fill=NAVY_500)
            textbox(slide, x + Inches(0.6), by, w - Inches(1.0),
                    Inches(0.4), b, font=FONT, size=11, color=INK_0,
                    line_spacing=1.35)
            by += Inches(0.42)

    mock_panel(x_left, "CLINICIAN COCKPIT", "For the treating oncologist",
               ["Overview", "Plan", "Trials", "Documents", "Clinical"],
               [
                   "Intake and pathology at a glance",
                   "Guideline pathway with reasoning at each step",
                   "Ranked Regeneron trials with clear verdicts",
                   "One-click clinician PDF report",
               ])
    mock_panel(x_right, "PATIENT SIDEBAR", "For the person it's about",
               ["Diagnosis", "Plan", "Healing", "Questions", "Next"],
               [
                   "Plain-English headlines, no jargon",
                   "The avatar reads answers aloud",
                   "Questions to ask the care team",
                   "Trial options explained in lay terms",
               ])


# SLIDE 9 — Readiness & roadmap
def slide_readiness(prs):
    slide = blank_slide(prs)
    chrome(slide, 8, "Readiness")
    headline(slide, "Open-source, reproducible, ready to run.", size=32)
    subhead(slide,
            "MIT licensed. One command for the backend, one for the frontend. "
            "If a key is missing, that feature switches off instead of crashing.",
            top=Inches(2.2))

    cols = [
        ("SHIPPING TODAY", NAVY_700, [
            "Pulls oncology fields out of PDFs",
            "Guideline walker with live reasoning",
            "Regeneron trial matching on a map",
            "One-click clinician PDF report",
            "Clinician and patient views",
        ]),
        ("IN PROGRESS", NAVY_500, [
            "Wire Kimi chat into the avatar",
            "Capture more intake (ECOG, RECIST, prior tx)",
            "Expand trial coverage past melanoma",
            "Pre-seeded demo case for judges",
        ]),
        ("PRIVACY & DEPLOY", INK_0, [
            "Synthetic data only. No real patient data.",
            "Cases live in memory, one session at a time",
            "Secrets in .env, a health endpoint shows setup",
            "Vercel for the frontend, Render or Fly for the backend",
        ]),
    ]
    y = Inches(3.4)
    h = Inches(3.45)
    gap = Inches(0.3)
    w = int((Inches(12.133) - gap * 2) / 3)
    x = Inches(0.6)
    for title, accent, items in cols:
        rect(slide, x, y, w, h, fill=WHITE, line=INK_800, line_width=Pt(0.75))
        rect(slide, x, y, Inches(0.06), h, fill=accent)
        textbox(slide, x + Inches(0.35), y + Inches(0.25),
                w - Inches(0.7), Inches(0.3),
                title, font=FONT, size=10, bold=True, color=accent, tracking=250)
        hline(slide, x + Inches(0.35), y + Inches(0.65),
              w - Inches(0.7), color=INK_800)
        by = y + Inches(0.8)
        for it in items:
            oval(slide, x + Inches(0.35), by + Inches(0.09), Inches(0.07),
                 fill=accent)
            textbox(slide, x + Inches(0.55), by, w - Inches(0.9),
                    Inches(0.55), it, font=FONT, size=11, color=INK_0,
                    line_spacing=1.35)
            by += Inches(0.52)
        x += w + gap


# SLIDE 10 — Links & thanks
def slide_thanks(prs):
    slide = blank_slide(prs)
    chrome(slide, 9, "Try it")
    rect(slide, Inches(0), Inches(0), Inches(0.22), Inches(7.5), fill=NAVY_700)

    textbox(slide, Inches(0.9), Inches(1.4), Inches(11.5), Inches(1.2),
            "Try Onkos.",
            font=FONT_LIGHT, size=72, color=INK_0, line_spacing=1.0)
    textbox(slide, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.8),
            "Clone it, add four API keys, upload a case.",
            font=FONT_LIGHT, size=22, color=INK_400, line_spacing=1.2)

    y = Inches(4.2)
    labels = [
        ("CODE",    "github.com/<your-org>/onkos"),
        ("DEMO",    "onkos.example.app"),
        ("VIDEO",   "youtu.be/<your-video-id>"),
        ("CONTACT", "team@onkos.app"),
    ]
    col_w = Inches(2.9)
    gap = Inches(0.15)
    x0 = Inches(0.9)
    for i, (label, val) in enumerate(labels):
        x = x0 + (col_w + gap) * i
        rect(slide, x, y, col_w, Inches(1.3), fill=WHITE,
             line=INK_800, line_width=Pt(0.75))
        textbox(slide, x + Inches(0.3), y + Inches(0.25),
                col_w - Inches(0.6), Inches(0.3),
                label, font=FONT, size=10, bold=True, color=NAVY_700, tracking=300)
        textbox(slide, x + Inches(0.3), y + Inches(0.65),
                col_w - Inches(0.6), Inches(0.55),
                val, font=FONT, size=12.5, color=INK_0, line_spacing=1.3)

    hline(slide, Inches(0.9), Inches(6.6), Inches(11.5), color=INK_800)
    textbox(slide, Inches(0.9), Inches(6.72), Inches(11.5), Inches(0.3),
            "THANK YOU  ·  TEAM ONKOS",
            font=FONT, size=10, bold=True, color=NAVY_700, tracking=400)


def build(out_path: Path):
    prs = new_presentation()
    slide_cover(prs)
    slide_problem(prs)
    slide_security(prs)
    slide_architecture(prs)
    slide_synthetic_data(prs)
    slide_regeneron(prs)
    slide_reasoning(prs)
    slide_design(prs)
    slide_readiness(prs)
    slide_thanks(prs)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"wrote {out_path}  ({out_path.stat().st_size / 1024:.1f} KB)")


if __name__ == "__main__":
    repo_root = Path(__file__).resolve().parent.parent
    build(repo_root / "onkos_pitch.pptx")
